from pptx import Presentation
from pptx.dml.color import RGBColor
#import re  # Import the 're' module for regular expressions (is not used for now)

sensitive_words = [
    # **Category 1: Company Name and Variations**
    "NewCo",  # <--- the *actual* company name**
    "Co",
    "Corp",
    "Inc",
    "Limited",
    "Ltd",
    "Group",
    "AG",
    "GmbH",
    "S.A.",

    # **Category 2: People Associated with the Company (if known - add more names if you have them)**
    #"John Doe",         # fill in the person name
    "CEO",    # Titles if names are not always used, or as additional search terms
    "CFO",
    "VP",
    "Executive",

    # **Category 3: Contact Information (Domains, generic email/phone patterns)**
    "https://newcosite.com/", # <--- **IMPORTANT: Replace with the *actual* client website domain (if known)!**
    "@newcosite.com", # <--- **IMPORTANT: Replace with the *actual* client website domain (if known)!**

    #to be used for regular expressions when added
    #".com",
    #"https://",
    #"http://",
    # "@", # Email domain
    #".com",
    #"info@", # Common generic emails
    #"sales@",
    #"office@",
    #"support@",

    # **Category 4: Project/Engagement Specific Terms (if applicable)**
    "Project", # Example project name
    "Phase",

    # **Category 5: Generic Sensitive Data Keywords (broad terms - might have false positives)**
    "Contact",
    "Confidential",
    "Acquisition",
    "Merger",
    "Proprietary",
    "Internal Use Only",
    "Sensitive Data",
    "Private",
    "Secret",
    "Financials",
    "Financial Data",
    "Revenue",
    "Profit",
    "Loss",
    "Budget",
    "Forecast",
    "Strategy",
    "Business Plan",
    "Competitors",
    "Competitor",
    "Client List",
    "Customer Data",
    "Personal Information",
    "Terms of Agreement",
    "Contract Details",
    "Legal Agreement",
    "Investment",
    "Valuation",
    "Share Price",
    "Salaries",
    "Compensation",
    "Employee Data",
    "Performance Reviews",
    "Location Data",
    "IP Address",
    "Social Security Number",
    "SSN",
    "TIN",

    # **Category 6: Industry-Specific Keywords**

]

presentation_path = "my_presentation.pptx"  # Replace with your actual file name
presentation = Presentation(presentation_path)
for slide in presentation.slides:
    for shape in slide.shapes:

        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:  # Iterate through runs within the paragraph
                    run_text = run.text  # Get text of the current run
                    is_sensitive_run = False
                    sensitive_word_found = None  # Variable to store the sensitive word that was found

                    for sensitive_word in sensitive_words:
                        if sensitive_word.lower() in run_text.lower():  # Case-insensitive check within the RUN's text
                            is_sensitive_run = True
                            sensitive_word_found = sensitive_word  # Store the sensitive word that caused the match
                            break

                    if is_sensitive_run:  # If the RUN is sensitive, modify it
                        run.font.color.rgb = RGBColor(255, 0, 0)  # Apply red color to the RUN
                        print(
                            f"Found sensitive word '{sensitive_word_found}' in run text: '{run_text}'")  # PRINT statement added here

        presentation.save("modified_presentation.pptx")
        #print(
                                #"Presentation processing and saving complete. Check 'modified_presentation.pptx'")  # Optional confirmation